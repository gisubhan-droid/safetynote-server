"""
DECK-2 공통편 — SafetyNOTE 사용자 설명서
대상: 전 권한 공통 (처음 사용하는 모든 사용자)
슬라이드: 8장
비율: 4:3 (25.4cm x 19.05cm)
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.dml.color import RGBColor
import copy

# ── 브랜드 색상 ─────────────────────────────────────────────
C_PINK      = RGBColor(0xD7, 0x00, 0x72)   # Primary #D70072
C_PURPLE    = RGBColor(0x68, 0x51, 0x82)   # Secondary #685182
C_DARK      = RGBColor(0x1F, 0x29, 0x37)   # 본문 텍스트
C_GRAY      = RGBColor(0x6B, 0x72, 0x80)   # 보조 텍스트
C_LGRAY     = RGBColor(0xF3, 0xF4, 0xF6)   # 배경 연회색
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER    = RGBColor(0xE5, 0xE7, 0xEB)
C_GREEN     = RGBColor(0x05, 0x96, 0x69)
C_ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)

# ── 슬라이드 크기: 4:3 ──────────────────────────────────────
W = Cm(25.4)
H = Cm(19.05)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃

# ════════════════════════════════════════════════════════════
# 헬퍼 함수
# ════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w if line_w else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_text_box(slide, text, x, y, w, h,
                 size=Pt(13), bold=False, color=C_DARK,
                 align=PP_ALIGN.LEFT, fill=None, line=None,
                 line_w=Pt(0), wrap=True):
    """배경색 있는 텍스트박스"""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w if line_w else Pt(1)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return shape

def slide_header(slide, title, subtitle=None, accent=C_PINK):
    """슬라이드 상단 헤더 바"""
    # 헤더 배경
    add_rect(slide, 0, 0, W, Cm(2.8), fill=accent)
    # 좌측 흰 악센트 바
    add_rect(slide, 0, 0, Cm(0.5), Cm(2.8), fill=C_WHITE)
    # 타이틀
    add_text(slide, title,
             Cm(1.0), Cm(0.35), Cm(20), Cm(1.4),
             size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Cm(1.0), Cm(1.75), Cm(22), Cm(0.8),
                 size=Pt(12), bold=False, color=RGBColor(0xFF,0xDD,0xF2),
                 align=PP_ALIGN.LEFT)

def slide_footer(slide, page_num, total=8, note=None):
    """슬라이드 하단 푸터"""
    add_rect(slide, 0, Cm(18.3), W, Cm(0.75), fill=C_LGRAY)
    add_text(slide, "Safety NOTE  |  사용자 설명서 — 공통편",
             Cm(0.5), Cm(18.35), Cm(18), Cm(0.6),
             size=Pt(9), color=C_GRAY, align=PP_ALIGN.LEFT)
    add_text(slide, f"{page_num} / {total}",
             Cm(22.5), Cm(18.35), Cm(2.5), Cm(0.6),
             size=Pt(9), color=C_GRAY, align=PP_ALIGN.RIGHT)
    if note:
        add_text(slide, f"💡 {note}",
                 Cm(0.5), Cm(17.5), Cm(24), Cm(0.7),
                 size=Pt(10), color=C_PURPLE, italic=True)

def screenshot_placeholder(slide, x, y, w, h, label="스크린샷"):
    """스크린샷 자리 표시자"""
    add_rect(slide, x, y, w, h, fill=RGBColor(0xF0,0xF4,0xFF),
             line=C_PURPLE, line_w=Pt(1.5))
    # 카메라 아이콘 텍스트
    add_text(slide, "📱",
             x, y + h*0.25, w, Cm(1.2),
             size=Pt(28), align=PP_ALIGN.CENTER)
    add_text(slide, label,
             x, y + h*0.55, w, Cm(0.8),
             size=Pt(11), color=C_PURPLE, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, "← 스크린샷 삽입",
             x, y + h*0.72, w, Cm(0.7),
             size=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

def step_badge(slide, x, y, number, color=C_PINK):
    """단계 번호 원형 배지"""
    add_rect(slide, x, y, Cm(0.9), Cm(0.9), fill=color)
    add_text(slide, str(number),
             x, y - Cm(0.02), Cm(0.9), Cm(0.94),
             size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def tip_box(slide, text, x, y, w, h=Cm(0.9)):
    """Tip / 주의 박스"""
    add_rect(slide, x, y, w, h,
             fill=RGBColor(0xFF,0xF7,0xED),
             line=C_ORANGE, line_w=Pt(1))
    add_text(slide, f"⚠  {text}",
             x + Cm(0.2), y + Cm(0.05), w - Cm(0.3), h,
             size=Pt(10.5), color=RGBColor(0x92,0x40,0x0E))

def good_box(slide, text, x, y, w, h=Cm(0.9)):
    """확인/완료 박스"""
    add_rect(slide, x, y, w, h,
             fill=RGBColor(0xEC,0xFD,0xF5),
             line=C_GREEN, line_w=Pt(1))
    add_text(slide, f"✓  {text}",
             x + Cm(0.2), y + Cm(0.05), w - Cm(0.3), h,
             size=Pt(10.5), color=RGBColor(0x06,0x5F,0x46))


# ════════════════════════════════════════════════════════════
# 슬라이드 01 — 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# 배경 그라디언트 효과 (두 영역으로 분할)
add_rect(s, 0, 0, W, H*0.55, fill=C_PINK)
add_rect(s, 0, H*0.55, W, H*0.45, fill=C_WHITE)

# 상단 장식 바
add_rect(s, 0, 0, Cm(0.7), H*0.55, fill=C_PURPLE)

# 앱 아이콘 자리 (흰 원형 배경)
add_rect(s, Cm(10.2), Cm(1.2), Cm(4.8), Cm(4.8),
         fill=C_WHITE, line=RGBColor(0xFF,0xCC,0xE8), line_w=Pt(2))
add_text(s, "🛡️", Cm(10.2), Cm(1.5), Cm(4.8), Cm(3.8),
         size=Pt(52), align=PP_ALIGN.CENTER)

# 앱 이름
add_text(s, "Safety NOTE",
         Cm(1.5), Cm(6.3), Cm(22), Cm(1.8),
         size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, "현장 안전관리 시스템",
         Cm(1.5), Cm(7.9), Cm(22), Cm(1.0),
         size=Pt(16), color=RGBColor(0xFF,0xCC,0xE8),
         align=PP_ALIGN.CENTER)

# 구분선
add_rect(s, Cm(3), Cm(10.55), Cm(19.4), Cm(0.06), fill=C_BORDER)

# 덱 제목
add_text(s, "사용자 설명서",
         Cm(1.5), Cm(10.8), Cm(22), Cm(1.3),
         size=Pt(26), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
add_text(s, "DECK-2  |  공통편 — 설치 · 로그인 · 화면 익히기",
         Cm(1.5), Cm(12.0), Cm(22), Cm(0.9),
         size=Pt(14), color=C_PURPLE, align=PP_ALIGN.CENTER)

# 대상 배지
add_text_box(s, "📋  대상: 처음 사용하는 모든 사용자",
             Cm(6.5), Cm(13.3), Cm(12.4), Cm(0.85),
             size=Pt(12), color=C_WHITE, fill=C_PURPLE,
             align=PP_ALIGN.CENTER, line_w=Pt(0))

# 하단 버전
add_text(s, "v2.0.0  |  2026",
         Cm(1.5), Cm(17.8), Cm(22), Cm(0.6),
         size=Pt(10), color=C_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 02 — 목차
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
slide_header(s, "이 설명서로 배울 내용", "목차 — 4단계로 쉽게 배웁니다")
slide_footer(s, 2)

items = [
    (1, "앱 설치",        "스마트폰에 Safety NOTE 앱 설치하기",           C_PINK),
    (2, "첫 로그인",      "아이디·비밀번호로 처음 로그인하기",              C_PURPLE),
    (3, "화면 구성 이해", "메뉴가 어디 있는지 한눈에 파악하기",            C_PINK),
    (4, "알림 확인",      "서명 요청 · 푸시 알림 받고 확인하기",            C_PURPLE),
]

for i, (num, title, desc, color) in enumerate(items):
    row_y = Cm(3.4) + i * Cm(3.1)

    # 카드 배경
    add_rect(s, Cm(1.2), row_y, Cm(22.5), Cm(2.6),
             fill=C_LGRAY, line=C_BORDER, line_w=Pt(1))

    # 번호 배지
    add_rect(s, Cm(1.2), row_y, Cm(2.0), Cm(2.6), fill=color)
    add_text(s, str(num),
             Cm(1.2), row_y + Cm(0.55), Cm(2.0), Cm(1.5),
             size=Pt(30), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 내용
    add_text(s, title,
             Cm(4.0), row_y + Cm(0.3), Cm(18), Cm(1.0),
             size=Pt(16), bold=True, color=C_DARK)
    add_text(s, desc,
             Cm(4.0), row_y + Cm(1.3), Cm(18), Cm(0.9),
             size=Pt(12), color=C_GRAY)

    # 화살표
    add_text(s, "›",
             Cm(22.5), row_y + Cm(0.7), Cm(1.2), Cm(1.2),
             size=Pt(22), color=color, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# 슬라이드 03 — APK 설치 (Android)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
slide_header(s, "STEP 1  앱 설치하기", "Android 스마트폰에 Safety NOTE 앱을 설치합니다")
slide_footer(s, 3, note="이미 설치되어 있다면 STEP 2(로그인)로 바로 이동하세요")

# 왼쪽: 단계 설명
steps = [
    ("로그인 화면 열기",
     "관리자에게 받은 주소를\n스마트폰 브라우저에 입력하세요"),
    ("APK 다운로드",
     "화면 하단 'APK 설치 파일\n다운로드' 버튼을 누르세요"),
    ("설치 허용",
     "'알 수 없는 출처 앱 허용'\n메시지가 나오면 허용을 누르세요"),
    ("설치 완료",
     "Safety NOTE 아이콘이 생기면\n설치가 완료된 것입니다 ✓"),
]

for i, (title, desc) in enumerate(steps):
    sy = Cm(3.2) + i * Cm(3.3)
    # 단계 배지
    step_badge(s, Cm(0.8), sy + Cm(0.05), i + 1)
    # 연결선
    if i < 3:
        add_rect(s, Cm(1.17), sy + Cm(0.95), Cm(0.06), Cm(2.4),
                 fill=RGBColor(0xE5,0xE7,0xEB))
    # 제목
    add_text(s, title,
             Cm(2.2), sy, Cm(12), Cm(0.85),
             size=Pt(14), bold=True, color=C_DARK)
    add_text(s, desc,
             Cm(2.2), sy + Cm(0.85), Cm(12), Cm(1.6),
             size=Pt(12), color=C_GRAY, wrap=True)

# 오른쪽: 스크린샷
screenshot_placeholder(s, Cm(14.8), Cm(3.0), Cm(9.8), Cm(13.5),
                        "로그인 화면\n(APK 다운로드 버튼 보임)")

tip_box(s, "PC나 태블릿은 설치 없이 브라우저에서 바로 접속 가능합니다",
        Cm(0.7), Cm(16.8), Cm(13.8))


# ════════════════════════════════════════════════════════════
# 슬라이드 04 — PC / 브라우저 접속
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
slide_header(s, "STEP 1-B  브라우저로 접속하기", "PC · 태블릿은 설치 없이 브라우저에서 바로 사용합니다")
slide_footer(s, 4)

# 두 칸 카드
cards = [
    ("🏢  사무실 · 현장 내부망",
     "사내 네트워크에 연결된 상태에서\n브라우저 주소창에 입력하세요",
     "https://NAS주소:3443",
     C_PINK),
    ("🌐  외부 / 집에서 접속",
     "인터넷이 연결된 어디서나\n아래 주소로 접속할 수 있습니다",
     "https://도메인주소:3443",
     C_PURPLE),
]

for i, (title, desc, addr, color) in enumerate(cards):
    cx = Cm(0.8) + i * Cm(12.4)
    # 카드
    add_rect(s, cx, Cm(3.2), Cm(11.8), Cm(8.5),
             fill=C_LGRAY, line=C_BORDER, line_w=Pt(1))
    # 상단 색상 바
    add_rect(s, cx, Cm(3.2), Cm(11.8), Cm(0.5), fill=color)
    add_text(s, title,
             cx + Cm(0.4), Cm(4.0), Cm(11), Cm(1.0),
             size=Pt(14), bold=True, color=C_DARK)
    add_text(s, desc,
             cx + Cm(0.4), Cm(5.1), Cm(11), Cm(1.6),
             size=Pt(12), color=C_GRAY, wrap=True)
    # 주소 박스
    add_rect(s, cx + Cm(0.4), Cm(7.0), Cm(11.0), Cm(1.1),
             fill=C_WHITE, line=color, line_w=Pt(1.5))
    add_text(s, addr,
             cx + Cm(0.4), Cm(7.05), Cm(11.0), Cm(1.0),
             size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)

# 공통 안내
add_rect(s, Cm(0.8), Cm(12.2), Cm(23.8), Cm(4.5),
         fill=RGBColor(0xF0,0xF4,0xFF), line=C_PURPLE, line_w=Pt(1))
add_text(s, "⚠  보안 경고 메시지가 나와도 당황하지 마세요!",
         Cm(1.2), Cm(12.4), Cm(22), Cm(0.8),
         size=Pt(13), bold=True, color=C_PURPLE)

warn_steps = [
    "Chrome :  '고급'  →  '안전하지 않은 사이트로 이동'  클릭",
    "Samsung 인터넷 :  '위험을 감수하고 계속'  클릭",
    "Safari :  '고급'  →  '안전하지 않은 웹 사이트 방문'  클릭",
]
for j, ws in enumerate(warn_steps):
    add_text(s, f"•  {ws}",
             Cm(1.4), Cm(13.3) + j * Cm(0.95), Cm(22.5), Cm(0.85),
             size=Pt(11.5), color=C_DARK)

good_box(s, "회사 내부 서버이므로 안전합니다. 그냥 '계속 진행'을 누르면 됩니다",
         Cm(0.8), Cm(16.6), Cm(23.8))


# ════════════════════════════════════════════════════════════
# 슬라이드 05 — 로그인
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
slide_header(s, "STEP 2  로그인하기", "관리자에게 받은 아이디와 비밀번호로 로그인합니다")
slide_footer(s, 5, note="첫 로그인 후에는 반드시 비밀번호를 변경하세요")

# 왼쪽: 절차
steps2 = [
    ("아이디 입력",   "관리자에게 받은 아이디를\n입력합니다 (전화번호 형식)"),
    ("비밀번호 입력", "관리자에게 받은\n임시 비밀번호를 입력합니다"),
    ("로그인 버튼",   "파란색 '로그인' 버튼을\n누르면 완료됩니다"),
]
for i, (title, desc) in enumerate(steps2):
    sy = Cm(3.3) + i * Cm(3.5)
    step_badge(s, Cm(0.8), sy + Cm(0.05), i + 1)
    if i < 2:
        add_rect(s, Cm(1.17), sy + Cm(0.95), Cm(0.06), Cm(2.6),
                 fill=RGBColor(0xE5,0xE7,0xEB))
    add_text(s, title,
             Cm(2.2), sy, Cm(12), Cm(0.85),
             size=Pt(14), bold=True, color=C_DARK)
    add_text(s, desc,
             Cm(2.2), sy + Cm(0.85), Cm(12), Cm(1.8),
             size=Pt(12), color=C_GRAY, wrap=True)

# 오른쪽: 스크린샷
screenshot_placeholder(s, Cm(14.8), Cm(3.0), Cm(9.8), Cm(13.5),
                        "로그인 화면")

# 첫 로그인 안내 박스
add_rect(s, Cm(0.7), Cm(13.8), Cm(13.5), Cm(2.8),
         fill=RGBColor(0xFF,0xF0,0xF7), line=C_PINK, line_w=Pt(1.5))
add_text(s, "🔑  첫 로그인 후 반드시 비밀번호 변경!",
         Cm(1.0), Cm(14.0), Cm(13), Cm(0.9),
         size=Pt(13), bold=True, color=C_PINK)
add_text(s, "메뉴 → 내 계정 → 비밀번호 변경\n임시 비밀번호는 보안에 취약합니다",
         Cm(1.0), Cm(14.9), Cm(13), Cm(1.5),
         size=Pt(11.5), color=C_DARK, wrap=True)


# ════════════════════════════════════════════════════════════
# 슬라이드 06 — 화면 구성 이해
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
slide_header(s, "STEP 3  화면 구성 이해하기", "로그인 후 보이는 화면을 한눈에 파악합니다")
slide_footer(s, 6)

# 오른쪽 스크린샷
screenshot_placeholder(s, Cm(14.5), Cm(3.0), Cm(10.2), Cm(13.8),
                        "메인 화면 전체\n(로그인 직후)")

# 왼쪽 설명
parts = [
    (Cm(0.7), Cm(3.2),  C_PINK,   "① 왼쪽 메뉴 레일",
     "화면 왼쪽에 아이콘들이 세로로\n줄지어 있습니다.\n아이콘을 누르면 해당 메뉴가\n오른쪽으로 펼쳐집니다."),
    (Cm(0.7), Cm(8.0),  C_PURPLE, "② 메인 내용 영역",
     "아이콘을 누른 후 원하는 메뉴를\n선택하면 이 영역에 내용이\n표시됩니다."),
    (Cm(0.7), Cm(12.0), C_GREEN,  "③ 상단 알림 아이콘",
     "오른쪽 위 🔔 아이콘을 누르면\n서명 요청·알림을 확인할 수\n있습니다."),
]

for px, py, color, title, desc in parts:
    # 좌측 색상 바
    add_rect(s, px, py, Cm(0.35), Cm(4.5), fill=color)
    add_text(s, title,
             px + Cm(0.6), py, Cm(13.0), Cm(0.85),
             size=Pt(14), bold=True, color=C_DARK)
    add_text(s, desc,
             px + Cm(0.6), py + Cm(0.9), Cm(13.0), Cm(3.2),
             size=Pt(12), color=C_GRAY, wrap=True)


# ════════════════════════════════════════════════════════════
# 슬라이드 07 — 알림 확인
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
slide_header(s, "STEP 4  알림 확인하기", "서명 요청이 오면 알림으로 알려드립니다")
slide_footer(s, 7)

# 오른쪽 스크린샷
screenshot_placeholder(s, Cm(14.5), Cm(3.0), Cm(10.2), Cm(13.8),
                        "알림 화면\n(벨 아이콘 클릭 후)")

# 왼쪽 알림 종류
add_text(s, "알림이 오는 경우",
         Cm(0.7), Cm(3.2), Cm(13.5), Cm(0.9),
         size=Pt(16), bold=True, color=C_DARK)

notif_items = [
    (C_PINK,   "서명 요청",   "감독자가 서명을 요청했습니다.\n빨간 배지 숫자 = 미처리 건수"),
    (C_PURPLE, "작업 배정",   "새 작업이 배정되었습니다.\n내 작업 목록에서 확인하세요"),
    (C_GREEN,  "공지 알림",   "관리자가 공지를 발송했습니다.\n알림 목록에서 내용을 확인하세요"),
]

for i, (color, title, desc) in enumerate(notif_items):
    ny = Cm(4.4) + i * Cm(3.4)
    add_rect(s, Cm(0.7), ny, Cm(13.3), Cm(3.0),
             fill=C_LGRAY, line=C_BORDER, line_w=Pt(1))
    add_rect(s, Cm(0.7), ny, Cm(0.4), Cm(3.0), fill=color)
    add_text(s, title,
             Cm(1.5), ny + Cm(0.2), Cm(12), Cm(0.85),
             size=Pt(14), bold=True, color=C_DARK)
    add_text(s, desc,
             Cm(1.5), ny + Cm(1.05), Cm(12), Cm(1.7),
             size=Pt(11.5), color=C_GRAY, wrap=True)

# 알림 확인 방법
add_rect(s, Cm(0.7), Cm(14.7), Cm(13.3), Cm(2.0),
         fill=RGBColor(0xF0,0xF4,0xFF), line=C_PURPLE, line_w=Pt(1))
add_text(s, "알림 확인 방법",
         Cm(1.0), Cm(14.9), Cm(12.7), Cm(0.75),
         size=Pt(13), bold=True, color=C_PURPLE)
add_text(s, "화면 오른쪽 위  🔔  아이콘 클릭  →  알림 목록 확인  →  해당 항목 클릭",
         Cm(1.0), Cm(15.65), Cm(12.7), Cm(0.85),
         size=Pt(11.5), color=C_DARK)


# ════════════════════════════════════════════════════════════
# 슬라이드 08 — 마무리 / 자주 묻는 질문
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=C_WHITE)
slide_header(s, "자주 묻는 질문 (FAQ)", "사용 중 불편한 점이 있을 때 먼저 확인해 보세요")
slide_footer(s, 8)

faqs = [
    ("Q.  로그인이 안 돼요",
     "아이디·비밀번호를 다시 확인하세요.\n그래도 안 되면 관리자에게 비밀번호 초기화를 요청하세요.",
     C_PINK),
    ("Q.  화면이 안 열려요 (인증서 오류)",
     "'고급'  또는  '세부정보'를 눌러 '계속 진행'을 선택하세요.\n회사 내부 서버라 정상입니다.",
     C_PURPLE),
    ("Q.  알림(벨)이 울리지 않아요",
     "스마트폰 설정 → 앱 → Safety NOTE → 알림 → 허용으로 바꿔주세요.",
     C_GREEN),
    ("Q.  앱이 최신 버전인지 확인하고 싶어요",
     "로그인 후 메뉴 → 업데이트 탭에서 현재 버전을 확인할 수 있습니다.",
     C_ORANGE),
]

for i, (q, a, color) in enumerate(faqs):
    fy = Cm(3.2) + i * Cm(3.3)
    add_rect(s, Cm(0.7), fy, Cm(23.9), Cm(2.9),
             fill=C_LGRAY, line=C_BORDER, line_w=Pt(1))
    add_rect(s, Cm(0.7), fy, Cm(0.4), Cm(2.9), fill=color)
    add_text(s, q,
             Cm(1.4), fy + Cm(0.2), Cm(22.8), Cm(0.85),
             size=Pt(13), bold=True, color=C_DARK)
    add_text(s, a,
             Cm(1.4), fy + Cm(1.05), Cm(22.8), Cm(1.7),
             size=Pt(11.5), color=C_GRAY, wrap=True)

# 마무리 문구
add_text_box(s, "추가 문의는 관리자에게 연락하세요",
             Cm(5.5), Cm(16.6), Cm(14.4), Cm(0.85),
             size=Pt(12), color=C_WHITE,
             fill=C_PURPLE, align=PP_ALIGN.CENTER, line_w=Pt(0))


# ════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════
OUT = "/home/user/webapp/docs/slides/DECK-2_공통편.pptx"
prs.save(OUT)
print(f"✅  저장 완료: {OUT}")
